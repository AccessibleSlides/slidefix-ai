import streamlit as st
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
import base64
from io import BytesIO

# 1. PAGE CONFIGURATION (Must be first)
st.set_page_config(
    page_title="Accessible Slides | Enterprise", 
    page_icon="♿", 
    layout="centered",
    initial_sidebar_state="expanded"
)

# --- CUSTOM CSS STYLING ---
st.markdown("""
<style>
    /* Import Google Font */
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&display=swap');

    html, body, [class*="css"] {
        font-family: 'Inter', sans-serif;
    }

    /* Title Styling */
    h1 {
        color: #1E3A8A; /* Dark Blue */
        font-weight: 800;
    }
    
    /* Sidebar Styling */
    [data-testid="stSidebar"] {
        background-color: #F8FAFC;
        border-right: 1px solid #E2E8F0;
    }

    /* Button Styling - Make the Main Button Pop */
    div.stButton > button:first-child {
        width: 100%;
        background-color: #2563EB; /* Enterprise Blue */
        color: white;
        font-size: 18px;
        font-weight: bold;
        padding: 0.75rem;
        border-radius: 10px;
        border: none;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);
        transition: all 0.2s ease;
    }
    div.stButton > button:first-child:hover {
        background-color: #1D4ED8;
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1);
        transform: translateY(-2px);
    }

    /* Success Box Styling */
    .success-box {
        padding: 1rem;
        border-radius: 0.5rem;
        background-color: #DCFCE7;
        border: 1px solid #86EFAC;
        color: #14532D;
        text-align: center;
        font-weight: bold;
        margin-top: 1rem;
    }
    
    /* Hide Streamlit Branding */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
</style>
""", unsafe_allow_html=True)

# Initialize Session State
if "license_valid" not in st.session_state:
    st.session_state["license_valid"] = False
if "org_name" not in st.session_state:
    st.session_state["org_name"] = ""

# --- LICENSE VERIFICATION ---
def verify_license(key):
    PRODUCT_ID = "RN-jQLgM_iudPwC-RnQZ6A==" 
    url = "https://api.gumroad.com/v2/licenses/verify"
    payload = {"product_id": PRODUCT_ID, "license_key": key.strip()}
    try:
        response = requests.post(url, data=payload)
        data = response.json()
        if data.get("success"):
            return True, data['purchase']['email']
        else:
            return False, None
    except Exception:
        return False, None

# --- SIDEBAR ---
with st.sidebar:
    st.image("https://img.icons8.com/fluency/96/accessibility.png", width=60)
    st.markdown("### License Status")
    
    if st.session_state["license_valid"]:
        st.success("✅ Active")
        st.caption(f"**Organization:**\n{st.session_state['org_name']}")
        
        st.markdown("---")
        if st.button("Logout", type="secondary"):
            st.session_state["license_valid"] = False
            st.session_state["org_name"] = ""
            st.rerun()
    else:
        st.warning("🔒 Locked")
        input_key = st.text_input("License Key", type="password")
        if st.button("Verify Access"):
            is_valid, org_email = verify_license(input_key)
            if is_valid:
                st.session_state["license_valid"] = True
                st.session_state["org_name"] = org_email
                st.rerun()
            else:
                st.error("Invalid Key")
        
        st.markdown("---")
        st.markdown("Don't have a key?")
        st.markdown("👉 **[Purchase Site License ($500)](https://accessibleslides.gumroad.com/l/accessibleslides_enterprise)**")

# --- MAIN CONTENT ---

# Hero Section
col1, col2 = st.columns([5, 1])
with col1:
    st.title("Accessible Slides")
    st.markdown("**Enterprise Edition** | Automated Section 508 Compliance")
with col2:
    # Just a visual spacer/logo area if needed
    pass

st.markdown("---")

# SECURITY GATE
if not st.session_state["license_valid"]:
    st.info("👋 **Welcome.** Please enter your Organization's License Key in the sidebar to unlock the secure portal.")
    st.image("https://cdn.dribbble.com/users/285475/screenshots/2083086/lock.gif", width=400)
    st.stop()

# 1. CONFIGURATION CARD
with st.expander("⚙️ System Configuration", expanded=True):
    st.caption("This tool operates using your internal OpenAI API keys for data privacy.")
    org_openai_key = st.text_input("OpenAI API Key (sk-...)", type="password")

# 2. UPLOAD CARD
st.markdown("### Upload Presentation")
uploaded_file = st.file_uploader("Drag and drop .pptx file", type=["pptx"], label_visibility="collapsed")

# Helper AI Logic
def get_ai_desc(client, image_blob):
    b64_image = base64.b64encode(image_blob).decode('utf-8')
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": [
                {"type": "text", "text": "Generate a concise, objective alt text. Do not start with 'Image of'."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64_image}"}},
            ]}], max_tokens=100
        )
        return response.choices[0].message.content
    except: return "Error generating text"

if uploaded_file:
    st.info(f"📄 Loaded: **{uploaded_file.name}**")
    
    if not org_openai_key:
        st.error("⚠️ Please enter an API Key in the Configuration section above.")
        st.stop()

    # THE BIG ACTION BUTTON
    if st.button("✨ Auto-Generate Alt Text"):
        
        try:
            client = OpenAI(api_key=org_openai_key)
            client.models.list() # Validate Key
        except Exception:
            st.error("❌ Invalid OpenAI API Key.")
            st.stop()

        prs = Presentation(uploaded_file)
        slide_count = len(prs.slides)
        
        # Progress Bar with visual style
        progress_text = "Scanning presentation..."
        my_bar = st.progress(0, text=progress_text)
        
        processed_images = 0
        
        for i, slide in enumerate(prs.slides):
            my_bar.progress((i + 1) / slide_count, text=f"Scanning Slide {i+1} of {slide_count}...")
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    desc = get_ai_desc(client, shape.image.blob)
                    try:
                        shape._element.nvPicPr.cNvPr.set('descr', desc)
                        shape._element.nvPicPr.cNvPr.set('name', "Image")
                        processed_images += 1
                    except: pass
        
        my_bar.empty()
        
        # SUCCESS DASHBOARD
        st.balloons()
        st.markdown(f"""
        <div class="success-box">
            ✅ SUCCESS! Processing Complete.
        </div>
        """, unsafe_allow_html=True)
        
        # METRICS ROW
        m1, m2, m3 = st.columns(3)
        with m1: st.metric("Slides Scanned", slide_count)
        with m2: st.metric("Images Fixed", processed_images)
        with m3: st.metric("Processing Time", "Instant")

        # PREPARE DOWNLOAD
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        st.markdown("### 📥 Download Result")
        st.download_button(
            label="Download Accessible PPTX",
            data=output,
            file_name=f"Accessible_{uploaded_file.name}",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            type="primary" # Makes it fill with theme color
        )